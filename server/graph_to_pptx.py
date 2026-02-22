"""
graph_to_pptx.py

ELK などでレイアウト済みの JSON グラフを、
PowerPoint のスライドに「角丸ノード＋中央線＋矢印コネクタ＋プロパティラベル」で
書き出すスクリプト。

- 外部からは:
    from graph_to_pptx import graph_json_to_pptx
    graph_json_to_pptx(data_dict, "out.pptx")

- CLI からは:
    python graph_to_pptx.py graph.json out.pptx
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Dict, List, Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor


# ======================
# データモデル
# ======================

@dataclass
class Node:
    id: str
    label: str
    x: float
    y: float
    width: float
    height: float


@dataclass
class Edge:
    id: str
    source: str
    target: str
    label: str | None = None


@dataclass
class Graph:
    nodes: List[Node]
    edges: List[Edge]


# ======================
# JSON → Graph 変換
# ======================

def parse_graph(data: Dict[str, Any]) -> Graph:
    """
    Canvas / ELK のレイアウト結果っぽい JSON から簡易 Graph を作る。

    期待する形式（例）:
    {
      "nodes": [
        {"id": "E22_1", "label": "E22 Human-Made Object | Vase 123",
         "x": 100, "y": 150, "width": 260, "height": 60},
        ...
      ],
      "edges": [
        {"id": "e1", "from": "E22_1", "to": "E42_1", "label": "P1 is identified by"},
        ...
      ]
    }
    """
    nodes: List[Node] = []
    for n in data.get("nodes", []):
        width = float(n.get("width", 180))
        height = float(n.get("height", 80))
        nodes.append(
            Node(
                id=str(n["id"]),
                label=str(n.get("label", "")),
                x=float(n.get("x", 0.0)),
                y=float(n.get("y", 0.0)),
                width=width,
                height=height,
            )
        )

    edges: List[Edge] = []
    for e in data.get("edges", []):
        src = e.get("from") or e.get("source")
        tgt = e.get("to") or e.get("target")
        edges.append(
            Edge(
                id=str(e.get("id", f"{src}_{tgt}")),
                source=str(src),
                target=str(tgt),
                label=e.get("label"),  # ← プロパティ名（P1...）が入る想定
            )
        )

    return Graph(nodes=nodes, edges=edges)


# ======================
# 座標変換
# ======================

class LayoutTransform:
    """
    ELK の座標系 → スライド座標系へのアフィン変換。

    ノード群のバウンディングボックスをスライドサイズの
    (1 - 2 * margin_ratio) くらいに収める。
    """

    def __init__(
        self,
        slide_width: int,
        slide_height: int,
        nodes: List[Node],
        margin_ratio: float = 0.05,
    ):
        self.slide_width = slide_width
        self.slide_height = slide_height

        xs = [n.x for n in nodes] or [0.0]
        ys = [n.y for n in nodes] or [0.0]

        self.min_x = min(xs)
        self.max_x = max(xs)
        self.min_y = min(ys)
        self.max_y = max(ys)

        self.margin_x = slide_width * margin_ratio
        self.margin_y = slide_height * margin_ratio

        span_x = (self.max_x - self.min_x) or 1.0
        span_y = (self.max_y - self.min_y) or 1.0

        usable_w = slide_width - 2 * self.margin_x
        usable_h = slide_height - 2 * self.margin_y

        scale_x = usable_w / span_x
        scale_y = usable_h / span_y

        # 縦横同じスケールを使う
        scale = min(scale_x, scale_y)

        # 要望: 全体をさらに 1/2 にする
        self.scale = scale * 0.5

    def x_to_emu(self, x: float) -> int:
        return int(self.margin_x + (x - self.min_x) * self.scale)

    def y_to_emu(self, y: float) -> int:
        return int(self.margin_y + (y - self.min_y) * self.scale)

    def size_to_emu(self, w: float, h: float) -> tuple[int, int]:
        # 幅はそのまま、高さだけほんの少し（1.15倍）大きくする
        return int(w * self.scale), int(h * self.scale * 1.15)


# ======================
# ノード描画（角丸＋中央線＋テキスト）
# ======================

def draw_nodes(slide, graph: Graph, tfm: LayoutTransform) -> Dict[str, Any]:
    """
    ノードをスライドに描き、id → shape の辞書を返す。

    - 角丸四角形（薄い単色背景＋黒枠＋黒字、影なし）
    - 中央に 3pt の水平線（矩形にコネクタ接続して一緒に動く・影なし）
    - ラベルは " | " を改行2つに変換して中央寄せ
    """
    shapes = slide.shapes
    id_to_shape: Dict[str, Any] = {}

    for node in graph.nodes:
        left = tfm.x_to_emu(node.x)
        top = tfm.y_to_emu(node.y)
        width, height = tfm.size_to_emu(node.width, node.height)

        # 角丸四角
        rect = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        id_to_shape[node.id] = rect

        # 塗り: 薄いグレー単色
        fill = rect.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)

        # 枠線: 黒
        rect.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
        rect.line.width = Pt(1.0)

        # 影なし
        rect.shadow.inherit = False
        rect.shadow.visible = False

        # " | " → 改行2つに変換して上下に分ける
        display_label = node.label.replace("|", "\n\n").strip()

        tf = rect.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = display_label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        # 中央の水平線（ストレートコネクタを横線として使い、
        # 四角形に対して左右の接続ポイントで begin/end_connect する）
        mid_y = top + height // 2
        line = shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            left,
            mid_y,
            left + width,
            mid_y,
        )
        # 0: 上中央, 1: 左中央, 2: 下中央, 3: 右中央
        line.begin_connect(rect, 1)  # 左辺中央
        line.end_connect(rect, 3)    # 右辺中央
        line.line.width = Pt(3.0)
        line.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
        # 影なし
        line.shadow.inherit = False
        line.shadow.visible = False

    return id_to_shape


# ======================
# コネクタ（矢印）描画
# ======================

def _add_arrowhead_to_target(connector) -> None:
    """
    コネクタの終点側 (tailEnd = to) に矢印ヘッドを付ける。

    begin_connect( from ) / end_connect( to ) で作っているので、
    tailEnd = to 側（欲しい側）に矢印が付く。
    """
    ln = connector.line._get_or_add_ln()
    ln.append(
        parse_xml(
            '<a:tailEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        )
    )


def _choose_connection_sites(src_shape, tgt_shape) -> tuple[int, int]:
    """
    2つの図形の重心の相対位置から、接続ポイント(0:上,1:左,2:下,3:右)を決める。
    """
    sx, sy, sw, sh = src_shape.left, src_shape.top, src_shape.width, src_shape.height
    tx, ty, tw, th = tgt_shape.left, tgt_shape.top, tgt_shape.width, tgt_shape.height

    src_cx = sx + sw // 2
    src_cy = sy + sh // 2
    tgt_cx = tx + tw // 2
    tgt_cy = ty + th // 2

    dx = tgt_cx - src_cx
    dy = tgt_cy - src_cy

    # 横方向の距離の方が大きければ左右を優先、そうでなければ上下
    if abs(dx) >= abs(dy):
        if dx >= 0:
            # 右向き: from 右 → to 左
            return 3, 1
        else:
            # 左向き: from 左 → to 右
            return 1, 3
    else:
        if dy >= 0:
            # 下向き: from 下 → to 上
            return 2, 0
        else:
            # 上向き: from 上 → to 下
            return 0, 2


def draw_edges(
    slide,
    graph: Graph,
    tfm: LayoutTransform,
    id_to_shape: Dict[str, Any],
) -> None:
    shapes = slide.shapes

    for edge in graph.edges:
        src_shape = id_to_shape.get(edge.source)
        tgt_shape = id_to_shape.get(edge.target)
        if src_shape is None or tgt_shape is None:
            continue  # 片方無ければスキップ

        # 初期位置は適当に中心点同士
        src_cx = src_shape.left + src_shape.width // 2
        src_cy = src_shape.top + src_shape.height // 2
        tgt_cx = tgt_shape.left + tgt_shape.width // 2
        tgt_cy = tgt_shape.top + tgt_shape.height // 2

        conn = shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            src_cx,
            src_cy,
            tgt_cx,
            tgt_cy,
        )

        # 図形に接続（ボックスを動かすと線も追従）
        from_site, to_site = _choose_connection_sites(src_shape, tgt_shape)
        conn.begin_connect(src_shape, from_site)
        conn.end_connect(tgt_shape, to_site)

        # to 側（tailEnd）に矢印を付ける
        _add_arrowhead_to_target(conn)
        conn.line.width = Pt(1.5)
        conn.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
        # 影なし
        conn.shadow.inherit = False
        conn.shadow.visible = False

        # エッジラベル（プロパティ名）を中点付近のテキストボックスとして描画し、
        # コネクタとグループ化する
        if edge.label:
            mid_x = (src_cx + tgt_cx) // 2
            mid_y = (src_cy + tgt_cy) // 2

            dx = abs(tgt_cx - src_cx)
            dy = abs(tgt_cy - src_cy)

            # 適当なサイズ（小さすぎないように下限を設ける）
            w = int(max(dx * 0.4, 150_000))  # 約 2cm 程度
            h = int(max(dy * 0.4, 80_000))   # 約 1cm 程度

            tb = shapes.add_textbox(mid_x - w // 2, mid_y - h // 2, w, h)
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = edge.label
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(15)  # 1.5倍くらい
            p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            tb.shadow.inherit = False
            tb.shadow.visible = False

            # コネクタ＋ラベルをグループ化（まとめて動かせるように）
            shapes.add_group_shape([conn, tb])


# ======================
# メイン API
# ======================

def graph_json_to_pptx(data: Dict[str, Any], output_path: str) -> None:
    """
    JSON グラフを受け取り、1 枚のスライドに
    角丸ノード＋中央線＋矢印コネクタ＋プロパティラベルを描いた PPTX を保存する。
    """
    graph = parse_graph(data)

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    tfm = LayoutTransform(slide_width, slide_height, graph.nodes)

    id_to_shape = draw_nodes(slide, graph, tfm)
    draw_edges(slide, graph, tfm, id_to_shape)

    prs.save(output_path)


# ======================
# CLI 用エントリポイント
# ======================

if __name__ == "__main__":
    import argparse
    import json
    from pathlib import Path

    parser = argparse.ArgumentParser(
        description="JSON グラフ → PowerPoint (.pptx) 変換ツール"
    )
    parser.add_argument("input_json", help="入力 JSON ファイルパス")
    parser.add_argument("output_pptx", help="出力 PPTX ファイルパス")

    args = parser.parse_args()

    in_path = Path(args.input_json)
    out_path = Path(args.output_pptx)

    with in_path.open(encoding="utf-8") as f:
        data = json.load(f)

    graph_json_to_pptx(data, str(out_path))
